from pathlib import Path
#!/usr/bin/env python3
"""
Insert 2 SOTA/Baseline slides into existing P5_Presentation.pptx.
Slides inserted AFTER slide 2 (problem statement).
Existing slides are NOT modified — only inserted.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy, os
from lxml import etree

PPTX_PATH = str(Path(__file__).parent / 'P5_Presentation.pptx')

# ── Warna tema (sama dengan file asli) ───────────────────────────────────────
C_NAVY   = RGBColor(0x0D, 0x2B, 0x55)
C_TEAL   = RGBColor(0x00, 0x7A, 0x8A)
C_ORANGE = RGBColor(0xE8, 0x6A, 0x1A)
C_AMBER  = RGBColor(0xF5, 0xA6, 0x23)
C_RED    = RGBColor(0xC0, 0x39, 0x2B)
C_GREEN  = RGBColor(0x1A, 0x7A, 0x4A)
C_DGRAY  = RGBColor(0x44, 0x44, 0x55)
C_LGRAY  = RGBColor(0xF4, 0xF6, 0xF8)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BGGRAY = RGBColor(0xF0, 0xF2, 0xF5)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation(PPTX_PATH)
blank_layout = prs.slide_layouts[6]

# ── Helper: insert slide at specific position ─────────────────────────────────
def move_slide_to(prs, from_idx, to_idx):
    """Move slide from from_idx to to_idx (0-based)."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    el = slides[from_idx]
    xml_slides.remove(el)
    xml_slides.insert(to_idx, el)

# ── Helper: add shape ─────────────────────────────────────────────────────────
def rect(sl, x, y, w, h, fill_rgb=None, line_rgb=None, line_w=Pt(0)):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    sh = sl.shapes.add_shape(1, x, y, w, h)
    if fill_rgb:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill_rgb
    else:
        sh.fill.background()
    if line_rgb:
        sh.line.color.rgb = line_rgb
        sh.line.width = line_w
    else:
        sh.line.fill.background()
    return sh

def txt(sl, text, x, y, w, h, size=11, bold=False, italic=False,
        color=C_NAVY, align=PP_ALIGN.LEFT, wrap=True):
    txb = sl.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = 'Calibri'
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return txb

def header_bar(sl, title, subtitle=None):
    h = Inches(1.15)
    rect(sl, 0, 0, SLIDE_W, h, fill_rgb=C_NAVY)
    rect(sl, 0, h - Pt(4), SLIDE_W, Pt(4), fill_rgb=C_TEAL)
    txt(sl, title, Inches(0.4), Inches(0.08), Inches(12.5), Inches(0.7),
        size=26, bold=True, color=C_WHITE)
    if subtitle:
        txt(sl, subtitle, Inches(0.4), Inches(0.72), Inches(12.5), Inches(0.35),
            size=13, color=C_AMBER)

def bg(sl):
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=C_BGGRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE A — Tinjauan Literatur & Gap Penelitian
# ════════════════════════════════════════════════════════════════════════════
slA = prs.slides.add_slide(blank_layout)
bg(slA)
header_bar(slA,
    'Tinjauan Literatur: Di Mana Posisi Penelitian Ini?',
    'State-of-the-art SDN load balancing — celah yang belum pernah diisi')

# ── Tabel SOTA (kiri) ────────────────────────────────────────────────────────
rect(slA, Inches(0.3), Inches(1.25), Inches(8.5), Inches(5.8),
     fill_rgb=C_WHITE)

# Header tabel
sota_header_cols = [
    (Inches(0.3),  Inches(2.6), 'Penelitian'),
    (Inches(2.95), Inches(1.5), 'Algoritma'),
    (Inches(4.5),  Inches(1.3), 'WC Dinamis?'),
    (Inches(5.85), Inches(3.0), 'Keterbatasan Utama'),
]
rect(slA, Inches(0.3), Inches(1.25), Inches(8.5), Inches(0.42),
     fill_rgb=C_NAVY)
for cx, cw, ch in sota_header_cols:
    txt(slA, ch, cx + Inches(0.06), Inches(1.28), cw - Inches(0.1),
        Inches(0.38), size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# Baris data SOTA
sota_rows = [
    ('Kaur et al. [1] 2025',       'SDN Survey',         '❌', 'Tidak ada eksperimen LB empiris'),
    ('Shona & Sharma [15] 2025',   'WRR (dinamis)',       '❌', 'Hanya WRR — tanpa IWRR/WLC'),
    ('Kumari et al. [14] 2024',    'Online RL',           '~',  'Distribusi per-flow tidak diukur'),
    ('Zhou et al. [31] 2024',      'Deep RL',             '~',  'Fokus energi/throughput agregat'),
    ('Baeturohman [16] 2024',      'WRR vs WLC (Nginx)', '❌', 'Bukan SDN/OpenFlow; statis'),
    ('Gaffke & Pukels. [13] 2008', 'Sainte-Laguë (teori)','N/A','Matematis; bukan SDN flow-level'),
    ('✅  Penelitian ini',          'WRR·IWRR·WLC·SL',   '✅ 3 skenario',
     'Pertama: konvergensi berkelanjutan pasca weight change'),
]

row_colors = [C_LGRAY, C_WHITE, C_LGRAY, C_WHITE, C_LGRAY, C_WHITE,
              RGBColor(0xE8, 0xF8, 0xF0)]

ry = Inches(1.67)
for ri, (paper, algo, wc, limit) in enumerate(sota_rows):
    rh = Inches(0.58) if ri == len(sota_rows)-1 else Inches(0.54)
    rect(slA, Inches(0.3), ry, Inches(8.5), rh, fill_rgb=row_colors[ri])
    is_last = ri == len(sota_rows) - 1
    fc = C_GREEN if is_last else C_DGRAY
    fc_algo = C_TEAL if is_last else C_DGRAY
    fc_wc   = (C_GREEN if wc == '✅ 3 skenario' else
               (C_RED if wc == '❌' else C_AMBER))

    txt(slA, paper, Inches(0.36), ry + Inches(0.05),
        Inches(2.5), rh, size=9, bold=is_last, color=fc)
    txt(slA, algo, Inches(2.95), ry + Inches(0.05),
        Inches(1.45), rh, size=9, bold=is_last, color=fc_algo,
        align=PP_ALIGN.CENTER)
    txt(slA, wc, Inches(4.5), ry + Inches(0.05),
        Inches(1.25), rh, size=9, bold=is_last, color=fc_wc,
        align=PP_ALIGN.CENTER)
    txt(slA, limit, Inches(5.85), ry + Inches(0.05),
        Inches(2.85), rh, size=9, bold=is_last, color=fc)
    ry += rh

# ── Panel kanan: Gap & Baseline ───────────────────────────────────────────────
rect(slA, Inches(9.05), Inches(1.25), Inches(4.05), Inches(2.75),
     fill_rgb=C_WHITE, line_rgb=C_RED, line_w=Pt(1.5))
txt(slA, '🔴  Gap Literatur',
    Inches(9.2), Inches(1.32), Inches(3.8), Inches(0.38),
    size=13, bold=True, color=C_RED)

gap_lines = [
    'Tidak ada penelitian yang:',
    '① Membandingkan WRR + IWRR + WLC',
    '    secara bersamaan',
    '② Dalam lingkungan SDN/OpenFlow',
    '③ Di bawah kondisi weight change',
    '    dinamis yang eksplisit',
    '④ Mengukur konvergensi berkelanjutan',
    '    sebagai metrik mandiri',
]
gy = Inches(1.72)
for gl in gap_lines:
    c = C_NAVY if gl.startswith('①②③④') or gl[0] in '①②③④' else C_DGRAY
    b = gl[0] in '①②③④'
    txt(slA, gl, Inches(9.2), gy, Inches(3.75), Inches(0.3),
        size=10, bold=b, color=c)
    gy += Inches(0.27)

# Baseline box
rect(slA, Inches(9.05), Inches(4.15), Inches(4.05), Inches(2.85),
     fill_rgb=RGBColor(0xE8, 0xF4, 0xFD),
     line_rgb=C_TEAL, line_w=Pt(1.5))
txt(slA, '🔵  Baseline Penelitian',
    Inches(9.2), Inches(4.22), Inches(3.8), Inches(0.38),
    size=13, bold=True, color=C_TEAL)

baselines = [
    ('WRR',  '[8] Katevenis 1991', 'Sequence-based, static'),
    ('IWRR', '[9] Shreedhar 1996', 'Interleaved, static'),
    ('WLC',  '[15][16] 2024–2025', 'Connection-aware'),
]
by = Inches(4.65)
for algo, ref, desc in baselines:
    rect(slA, Inches(9.15), by, Inches(0.55), Inches(0.3),
         fill_rgb=C_NAVY)
    txt(slA, algo, Inches(9.15), by, Inches(0.55), Inches(0.3),
        size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slA, ref, Inches(9.75), by, Inches(1.5), Inches(0.3),
        size=9, bold=True, color=C_NAVY)
    txt(slA, desc, Inches(11.3), by, Inches(1.65), Inches(0.3),
        size=9, color=C_DGRAY)
    by += Inches(0.42)

txt(slA, '→ Sainte-Laguë [11][12][13] sebagai\n   kandidat solusi dari teori alokasi',
    Inches(9.2), Inches(6.0), Inches(3.75), Inches(0.7),
    size=10, bold=True, color=C_ORANGE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE B — Narasi Kutipan: Mengapa Ketiga Algoritma Ini?
# ════════════════════════════════════════════════════════════════════════════
slB = prs.slides.add_slide(blank_layout)
bg(slB)
header_bar(slB,
    'Baseline: Mengapa WRR, IWRR, dan WLC?',
    'Tiga kategori berbeda scheduler berbobot yang paling banyak digunakan di produksi SDN')

# ── 3 kolom baseline ──────────────────────────────────────────────────────────
algo_cards = [
    ('WRR',
     'Weighted Round Robin',
     C_NAVY,
     '[8] Katevenis et al. 1991\nIEEE J. Sel. Areas Commun.',
     'Membangun urutan siklik L = Σwᵢ.\n'
     'Proporsional sempurna di batas siklus,\n'
     'berosilasi di antaranya.\n\n'
     'Paling banyak digunakan di\n'
     'Linux LVS, HAProxy, Nginx.',
     'Structural Lock:\nRebuild + reset total saat WC',
     C_RED),

    ('IWRR',
     'Interleaved WRR',
     C_TEAL,
     '[9] Shreedhar & Varghese 1996\nIEEE/ACM Trans. Netw. (DRR)',
     'Interleaving lebih merata dalam siklus.\n'
     'Fairness lebih baik dari WRR untuk\n'
     'paket panjang bervariasi.\n\n'
     'Respons weight change:\n'
     'identik dengan WRR.',
     'Sequence Reset:\nRebuild + reset — sama dengan WRR',
     C_ORANGE),

    ('WLC',
     'Weighted Least Connections',
     RGBColor(0x6A, 0x0D, 0x0D),
     '[15] Shona & Sharma 2025  ETASR\n[16] Baeturohman & Santoso 2024',
     'Pilih server: argmin(active(i)/w(i)).\n'
     'Tidak ada sequence rebuild.\n'
     'Default scheduler Linux Virtual\n'
     'Server (LVS).',
     'Greedy Overshoot:\nHard-clearing → 100% ke srv1',
     C_RED),
]

for ci, (abbr, name, col, ref, desc, failure, fcol) in enumerate(algo_cards):
    bx = Inches(0.3) + ci * Inches(4.3)
    bw = Inches(4.1)

    # Card background
    rect(slB, bx, Inches(1.25), bw, Inches(5.8),
         fill_rgb=C_WHITE, line_rgb=col, line_w=Pt(1.5))
    rect(slB, bx, Inches(1.25), bw, Pt(5), fill_rgb=col)

    # Abbr badge
    rect(slB, bx + Inches(0.15), Inches(1.35), Inches(0.85), Inches(0.48),
         fill_rgb=col)
    txt(slB, abbr, bx + Inches(0.15), Inches(1.35), Inches(0.85), Inches(0.48),
        size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Algorithm name
    txt(slB, name, bx + Inches(1.1), Inches(1.38), bw - Inches(1.2),
        Inches(0.45), size=13, bold=True, color=col)

    # Reference
    rect(slB, bx + Inches(0.15), Inches(1.9), bw - Inches(0.3), Inches(0.55),
         fill_rgb=C_LGRAY)
    txt(slB, ref, bx + Inches(0.22), Inches(1.93),
        bw - Inches(0.4), Inches(0.5),
        size=9, italic=True, color=C_NAVY)

    # Description
    txt(slB, desc, bx + Inches(0.18), Inches(2.55),
        bw - Inches(0.35), Inches(2.4),
        size=10, color=C_DGRAY)

    # Failure mode
    rect(slB, bx + Inches(0.15), Inches(5.3), bw - Inches(0.3), Inches(1.5),
         fill_rgb=RGBColor(0xFF, 0xEB, 0xEB),
         line_rgb=fcol, line_w=Pt(1))
    txt(slB, '⚠  Kegagalan saat Weight Change:',
        bx + Inches(0.22), Inches(5.37), bw - Inches(0.4), Inches(0.3),
        size=9, bold=True, color=fcol)
    txt(slB, failure, bx + Inches(0.22), Inches(5.68),
        bw - Inches(0.4), Inches(0.9),
        size=11, bold=True, color=fcol)

# Bottom note
rect(slB, Inches(0.3), Inches(7.1), Inches(12.75), Inches(0.28),
     fill_rgb=C_NAVY)
txt(slB,
    'Ketiga algoritma ini dipilih karena merepresentasikan tiga kategori berbeda weighted scheduler yang paling banyak digunakan — '
    'masing-masing memiliki mekanisme failure yang berbeda secara fundamental saat weight change terjadi.',
    Inches(0.45), Inches(7.12), Inches(12.5), Inches(0.25),
    size=9, color=C_WHITE)


# ════════════════════════════════════════════════════════════════════════════
# Insert slides at position 3 (after slide 2 = problem slide)
# Slide A → position index 2 (0-based), Slide B → position index 3
# ════════════════════════════════════════════════════════════════════════════
total = len(prs.slides)
# Slide A was added at index total-2, Slide B at total-1
# Move Slide A to index 2 (after slide 1=title, slide 2=problem)
move_slide_to(prs, total - 2, 2)
# Now Slide B is at total-1, move to index 3
move_slide_to(prs, total - 1, 3)

# Save
prs.save(PPTX_PATH)
print(f'Saved: {PPTX_PATH}')
print(f'Total slides: {len(prs.slides)}')
for i, sl in enumerate(prs.slides, 1):
    texts = []
    for sh in sl.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                t = para.text.strip()
                if t and len(t) > 5:
                    texts.append(t[:65])
                    break
        if texts: break
    label = texts[0] if texts else '(no text)'
    print(f'  Slide {i:2d}: {label}')
